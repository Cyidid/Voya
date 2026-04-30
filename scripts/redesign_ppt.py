#!/usr/bin/env python3
"""
Voya PPT 重设计脚本 v2
- 字体：英文 Arial / 中文 宋体
- 封面：深色背景 + 云游图标 + ✈ 装饰 + 铜色品牌排版
- 章节分隔页：深色处理，铜色编号
- 内容页：底部品牌栏 + 铜色标题线
- 内容更新：修正过时信息
"""

import os
import re
from lxml import etree
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ─── 文件路径 ────────────────────────────────────────────────────────────────
SRC       = '/Users/cydid/Desktop/Voya · AI 旅行规划系统 (1).pptx'
DST       = '/Users/cydid/Desktop/Voya_Redesigned.pptx'
ICON_PATH = '/Users/cydid/Desktop/projects/assets/voya-icon.png'

# ─── 品牌色板 ────────────────────────────────────────────────────────────────
C_DARK    = RGBColor(0x12, 0x0F, 0x0B)   # 深棕黑封面背景
C_DARK2   = RGBColor(0x1E, 0x18, 0x12)   # 章节页深色
C_COPPER  = RGBColor(0xC4, 0x85, 0x4A)   # 铜色主题色
C_AMBER   = RGBColor(0xD4, 0xA2, 0x64)   # 浅铜/琥珀
C_CREAM   = RGBColor(0xFA, 0xF5, 0xEE)   # 暖米色（内容页背景）
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT    = RGBColor(0x1F, 0x23, 0x29)   # 正文深色
C_MUTED   = RGBColor(0x6B, 0x72, 0x80)   # 辅助灰色
C_BORDER  = RGBColor(0xE5, 0xDD, 0xCF)   # 分隔线

W_EMU = int(13.33 * 914400)
H_EMU = int(7.50 * 914400)

# ─── 辅助函数 ────────────────────────────────────────────────────────────────

def has_cjk(text: str) -> bool:
    return any('一' <= c <= '鿿' or '　' <= c <= '〿' for c in text)

def apply_run_font(run):
    """为每个 run 设置字体：英文 Arial / 中文 宋体"""
    t = run.text
    if not t.strip():
        return
    rPr = run._r.get_or_add_rPr()
    # 清理旧的字体声明
    for tag in [qn('a:latin'), qn('a:ea'), qn('a:cs')]:
        for el in rPr.findall(tag):
            rPr.remove(el)
    # 写入新字体
    latin = etree.SubElement(rPr, qn('a:latin'))
    ea    = etree.SubElement(rPr, qn('a:ea'))
    cs    = etree.SubElement(rPr, qn('a:cs'))
    latin.set('typeface', 'Arial')
    if has_cjk(t):
        ea.set('typeface', '宋体')
        cs.set('typeface', '宋体')
    else:
        ea.set('typeface', 'Arial')
        cs.set('typeface', 'Arial')

def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_shape_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def set_text_color(shape, color: RGBColor):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = color

def add_rect(slide, left_in, top_in, w_in, h_in, color: RGBColor, no_line=True):
    from pptx.util import Inches
    box = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left_in), Inches(top_in),
        Inches(w_in), Inches(h_in)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    if no_line:
        box.line.fill.background()
    return box

def add_textbox(slide, left_in, top_in, w_in, h_in,
                text, size_pt, color: RGBColor,
                bold=False, align=PP_ALIGN.LEFT, italic=False):
    from pptx.util import Inches, Pt
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in),
        Inches(w_in), Inches(h_in)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    apply_run_font(run)
    tb.line.fill.background()
    return tb


# ─── 加载原始 PPT ──────────────────────────────────────────────────────────
prs = Presentation(SRC)
total = len(prs.slides)
print(f'已加载：{total} 页幻灯片')

# ══════════════════════════════════════════════════════════════════════════════
# PASS 1：全局字体替换  Arial (EN) + 宋体 (ZH)
# ══════════════════════════════════════════════════════════════════════════════
print('Pass 1: 替换全局字体...')
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    apply_run_font(run)

# ══════════════════════════════════════════════════════════════════════════════
# PASS 2：修正过时内容
# ══════════════════════════════════════════════════════════════════════════════
print('Pass 2: 修正文本内容...')

FIXES = {
    '南亚（3）': '中东&海岛（3）',
    '南亚(3)': '中东&海岛（3）',
    '包含7个城市分布均匀': '包含25个城市（7个城市特征维度用于分类）',
    'Qwen LLM + Tavily 实时信息搜索': 'Qwen（通义千问）+ Tavily 实时搜索',
    '通义千问 (Qwen)，提供': '通义千问（Qwen），提供',
    'Cu Enpu': 'Gu Enpu',          # 修正拼写
}

def fix_text_in_slide(slide, fixes):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in fixes.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)

for slide in prs.slides:
    fix_text_in_slide(slide, FIXES)

# ══════════════════════════════════════════════════════════════════════════════
# PASS 3：封面重设计（第1页 index=0）
# ══════════════════════════════════════════════════════════════════════════════
print('Pass 3: 重设计封面...')

cover = prs.slides[0]

# 3-a: 深色底 (覆盖背景色)
set_slide_bg(cover, C_DARK)

# 3-b: 把原来的 FAF5EE 纯色填充矩形改成深色
for shape in cover.shapes:
    try:
        if shape.fill.type and str(shape.fill.type) == 'SOLID (1)':
            col_hex = str(shape.fill.fore_color.rgb).upper()
            if col_hex in ('FAF5EE', 'F5F3EF', 'FFFFFF'):
                set_shape_fill(shape, C_DARK)
    except Exception:
        pass

# 3-c: 铜色顶部装饰线（细线）
add_rect(cover, 0, 0, 13.33, 0.06, C_COPPER)

# 3-d: 左侧铜色竖线装饰
add_rect(cover, 0.72, 1.1, 0.05, 4.8, C_COPPER)

# 3-e: 大型装饰 ✈（右侧淡铜色，作为背景感）
tb_plane = cover.shapes.add_textbox(
    Inches(7.5), Inches(0.8),
    Inches(5.5), Inches(5.5)
)
tf = tb_plane.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.RIGHT
run = p.add_run()
run.text = '✈'
run.font.size = Pt(260)
run.font.color.rgb = RGBColor(0x2C, 0x1E, 0x12)   # 极深铜，近乎融入背景
tb_plane.line.fill.background()
# 发送到最底层
spTree = cover.shapes._spTree
spTree.remove(tb_plane._element)
spTree.insert(2, tb_plane._element)

# 3-f: 修改现有文字颜色（按 shape 整体判断，避免 run 级别遗漏）
for shape in cover.shapes:
    if not shape.has_text_frame:
        continue
    shape_text = shape.text_frame.text.strip()
    # 判断该 shape 属于哪类
    is_title    = 'Voya' in shape_text and 'AI 旅行' in shape_text
    is_subtitle = '负责任' in shape_text or '范式' in shape_text or '旅行规划对比' in shape_text
    is_author   = 'Group' in shape_text or 'Wu' in shape_text or 'Chen' in shape_text
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            t = run.text.strip()
            if not t:
                continue
            if is_title:
                run.font.color.rgb = C_COPPER
                run.font.bold = True
                run.font.size = Pt(42)
            elif is_subtitle:
                run.font.color.rgb = C_AMBER   # 整个副标题 shape 全部用琥珀色
                run.font.size = Pt(20)
            elif is_author:
                run.font.color.rgb = C_MUTED
                run.font.size = Pt(14)
            apply_run_font(run)

# 3-g: 添加 Voya 图标（右下角）
if os.path.exists(ICON_PATH):
    cover.shapes.add_picture(
        ICON_PATH,
        Inches(11.5), Inches(5.8),
        Inches(1.5), Inches(1.5)
    )
    print('  ✓ Voya 图标已添加至封面')

# 3-h: 底部深色品牌条
add_rect(cover, 0, 7.1, 13.33, 0.4, RGBColor(0x0A, 0x08, 0x06))
add_textbox(cover, 0.72, 7.15, 8.0, 0.3,
            'Responsible AI · Three Paradigms · Travel Planning System',
            9, RGBColor(0x6B, 0x52, 0x38),
            bold=False, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# PASS 4：章节分隔页深色处理
# 章节分隔页（0-indexed）：2, 6, 11, 17, 24, 28, 36, 40
# ══════════════════════════════════════════════════════════════════════════════
print('Pass 4: 处理章节分隔页...')

SECTION_SLIDES = [2, 6, 11, 17, 24, 28, 36, 40]

for idx in SECTION_SLIDES:
    slide = prs.slides[idx]
    set_slide_bg(slide, C_DARK2)

    for shape in slide.shapes:
        # 把原始浅色背景矩形改为深色
        try:
            if shape.fill.type and str(shape.fill.type) == 'SOLID (1)':
                col_hex = str(shape.fill.fore_color.rgb).upper()
                if col_hex in ('FAF5EE', 'F5F3EF', 'FFFFFF', 'F9F5F0'):
                    set_shape_fill(shape, C_DARK2)
        except Exception:
            pass

        if not shape.has_text_frame:
            continue

        text_all = shape.text_frame.text.strip()
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                t = run.text.strip()
                if not t:
                    continue
                # 大数字 → 铜色大字
                if re.fullmatch(r'0?\d', t):
                    run.font.color.rgb = C_COPPER
                    run.font.size = Pt(72)
                    run.font.bold = True
                    apply_run_font(run)
                # 主标题（中文）
                elif has_cjk(t) and len(t) > 2 and not t.startswith('（'):
                    run.font.color.rgb = C_WHITE
                    run.font.bold = True
                    apply_run_font(run)
                # 英文副标题
                elif t and not has_cjk(t) and len(t) > 3:
                    run.font.color.rgb = C_AMBER
                    run.font.size = Pt(14)
                    apply_run_font(run)
                # 括号说明
                else:
                    run.font.color.rgb = RGBColor(0xA0, 0x90, 0x80)
                    apply_run_font(run)

# ══════════════════════════════════════════════════════════════════════════════
# PASS 5：最后一页（Q&A，index=43）深色处理
# ══════════════════════════════════════════════════════════════════════════════
print('Pass 5: 处理结尾页...')

last = prs.slides[43]
set_slide_bg(last, C_DARK)
for shape in last.shapes:
    try:
        if shape.fill.type and str(shape.fill.type) == 'SOLID (1)':
            col_hex = str(shape.fill.fore_color.rgb).upper()
            if col_hex in ('FAF5EE', 'F5F3EF', 'FFFFFF'):
                set_shape_fill(shape, C_DARK)
    except Exception:
        pass
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            t = run.text.strip()
            if 'Q' in t and 'A' in t:
                run.font.color.rgb = C_COPPER
                run.font.size = Pt(60)
                run.font.bold = True
            elif t:
                run.font.color.rgb = C_WHITE
            apply_run_font(run)

# 加图标到结尾页
if os.path.exists(ICON_PATH):
    last.shapes.add_picture(
        ICON_PATH,
        Inches(6.0), Inches(2.5),
        Inches(1.3), Inches(1.3)
    )

# ══════════════════════════════════════════════════════════════════════════════
# PASS 6：内容页 — 添加底部品牌栏 + 标题铜色线（跳过封面/章节/结尾）
# ══════════════════════════════════════════════════════════════════════════════
print('Pass 6: 添加底部品牌栏...')

SKIP_SLIDES = set([0, 43] + SECTION_SLIDES)

for i, slide in enumerate(prs.slides):
    if i in SKIP_SLIDES:
        continue

    # 薄铜线（标题下方），判断标题 shape 位置
    # 所有内容页：顶部 0.56" 附近有标题行
    add_rect(slide, 0.83, 1.28, 11.67, 0.025, C_COPPER)

    # 底部品牌条（7.15" → 7.50"）
    add_rect(slide, 0, 7.12, 13.33, 0.38, RGBColor(0x1E, 0x18, 0x12))

    # 左侧 Voya 文字
    add_textbox(slide, 0.5, 7.17, 3.0, 0.3,
                'Voya · AI Travel Planner',
                8.5, RGBColor(0x80, 0x68, 0x48),
                bold=False)

    # 右侧页码
    page_num = i + 1
    add_textbox(slide, 11.8, 7.17, 1.3, 0.3,
                str(page_num),
                8.5, RGBColor(0x80, 0x68, 0x48),
                bold=False, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# PASS 6b：内容页浅色背景上的白色文字修复
# ══════════════════════════════════════════════════════════════════════════════
print('Pass 6b: 修复内容页白色字...')

LIGHT_BG_SLIDES = set(range(total)) - set([0, 43] + SECTION_SLIDES)

for i in LIGHT_BG_SLIDES:
    slide = prs.slides[i]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            fill_hex = str(shape.fill.fore_color.rgb).upper() if str(shape.fill.type) == 'SOLID (1)' else ''
            shape_is_dark = fill_hex in ('C4854A','1E1812','120F0B','0A0806','2C241C')
        except:
            shape_is_dark = False

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                try:
                    col_hex = str(run.font.color.rgb).upper()
                    if col_hex in ('FFFFFF','FAFAFA','F5F5F5','F0F0F0') and not shape_is_dark:
                        run.font.color.rgb = C_TEXT
                except:
                    pass

# ══════════════════════════════════════════════════════════════════════════════
# PASS 7：目录页（index=1）微调
# ══════════════════════════════════════════════════════════════════════════════
print('Pass 7: 优化目录页...')
toc = prs.slides[1]
for shape in toc.shapes:
    if not shape.has_text_frame:
        continue
    text_all = shape.text_frame.text.strip()
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            t = run.text.strip()
            if not t:
                continue
            # 数字标号 → 铜色
            if re.fullmatch(r'0\d', t):
                run.font.color.rgb = C_COPPER
                run.font.bold = True
                apply_run_font(run)
            # 章节主题 → 深色加粗
            elif has_cjk(t) and len(t) > 4:
                run.font.color.rgb = C_TEXT
                run.font.bold = True
                apply_run_font(run)
            # 小字描述 → 灰色
            elif t and len(t) > 5:
                run.font.color.rgb = C_MUTED
                apply_run_font(run)

# ══════════════════════════════════════════════════════════════════════════════
# 保存
# ══════════════════════════════════════════════════════════════════════════════
print(f'正在保存到 {DST} ...')
prs.save(DST)
print('✓ 完成！')
print(f'  输出：{DST}')
